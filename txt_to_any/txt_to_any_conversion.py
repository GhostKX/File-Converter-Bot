from user_states import UserFlow
import buttons as bt
import common_funtions as cf

from aiogram import F, Router
from aiogram.types import Message
from aiogram.fsm.context import FSMContext

import os
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

from pptx import Presentation
from pptx.util import Pt
import textwrap
import pandas

txt_router = Router()

################################### ***** TXT converting to DOCX, PDF, PPT, XLSX, CSV ***** ############################


@txt_router.message(F.text, UserFlow.txt_converting_file_type)
async def txt_file_converting_to(message: Message, state: FSMContext):
    user_id = message.from_user.id
    sending_key = f'txt_file_{user_id}'
    file_path = f'./txt_to_any/txt_files/'
    txt_file_path = await cf.getting_initial_file_path(state, sending_key, file_path, file_type='txt')

    converted_file_path = f'txt_to_any/txt_converted_files'

    ####################################################################################################################
                                                # TXT File converting to DOCX

    if message.text == '📝 To Word Docx':

        converted_file_path = await cf.creating_unique_file_name(user_id, converted_file_path, file_type='docx')

        if os.path.exists(txt_file_path):
            word_document = Document()

            with open(txt_file_path, 'r', encoding='utf-8') as file:
                for line in file:
                    word_document.add_paragraph(line.strip())

            word_document.save(converted_file_path)

            await cf.sending_converted_file(message, state, user_id, txt_file_path, converted_file_path)

    ####################################################################################################################
                                                # TXT File converting to PDF

    elif message.text == '📖 To PDF Document':

        converted_file_path = await cf.creating_unique_file_name(user_id, converted_file_path, file_type='pdf')

        if os.path.exists(txt_file_path):
            pdf_file = canvas.Canvas(converted_file_path, pagesize=letter)
            width, height = letter
            print(height)

            pdf_file.setFont('Times-Roman', 12)
            margin = 40

            y_position = height - margin
            print(y_position)

            with open(txt_file_path, 'r', encoding='UTF-8') as file:
                for line in file:
                    if y_position < margin:
                        pdf_file.showPage()
                        y_position = height - y_position

                    pdf_file.drawString(margin, y_position, line.strip())
                    y_position -= 15

            pdf_file.save()

        await cf.sending_converted_file(message, state, user_id, txt_file_path, converted_file_path)

    ####################################################################################################################
                                                # TXT File converting to PPT

    elif message.text == '📊 To PPT Presentation':
        converted_file_path = await cf.creating_unique_file_name(user_id, converted_file_path, file_type='pptx')
        template_file_path = f'./txt_to_any/txt_to_pptx_core_file/Core_file.pptx'

        if os.path.exists(txt_file_path) and os.path.exists(template_file_path):
            presentation = Presentation(template_file_path)
            slide_template = presentation.slides[0]

            text_box = slide_template.shapes[0]
            text_frame = text_box.text_frame
            text_box_width = text_box.width
            text_box_height = text_box.height
            left = text_box.left
            top = text_box.top
            max_lines_number = 21

            with open(txt_file_path, 'r', encoding='UTF-8') as txt_file:
                all_text = txt_file.readlines()

            line_count = 0

            for line in all_text:
                wrapped_lines = textwrap.wrap(line.strip(), width=124)

                for wrapped_line in wrapped_lines:
                    if line_count >= max_lines_number:
                        current_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                        new_text_box = current_slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
                        text_frame = new_text_box.text_frame
                        text_frame.word_wrap = True
                        line_count = 0

                    p = text_frame.add_paragraph()
                    p.text = wrapped_line
                    p.font_name = 'Arial'
                    p.font_size = Pt(16)
                    line_count += 1

            presentation.save(converted_file_path)

            await cf.sending_converted_file(message, state, user_id, txt_file_path, converted_file_path)

    ####################################################################################################################
                                                # TXT File converting to XLSX

    elif message.text == '📈 To Excel Spreadsheet':
        converted_file_path = await cf.creating_unique_file_name(user_id, converted_file_path, file_type='xlsx')

        if os.path.exists(txt_file_path):
            with open(txt_file_path, 'r', encoding='utf-8') as file:
                lines = file.readlines()

            df = pandas.DataFrame({'Content': [line.strip() for line in lines]})
            df.to_excel(converted_file_path, index=False)

            await cf.sending_converted_file(message, state, user_id, txt_file_path, converted_file_path)

    ####################################################################################################################
                                                # TXT File converting to CSV

    elif message.text == '🧾 To CSV File':
        converted_file_path = await cf.creating_unique_file_name(user_id, converted_file_path, file_type='csv')
        if os.path.exists(txt_file_path):
            with open(txt_file_path, 'r', encoding='utf-8') as file:
                lines = file.readlines()

            df = pandas.DataFrame({'Content': [line.strip() for line in lines]})
            df.to_csv(converted_file_path, index=False)

            await cf.sending_converted_file(message, state, user_id, txt_file_path, converted_file_path)

    ####################################################################################################################
                                                # Handling Back button logic

    elif message.text == '⏮️ Back':
        await cf.removing_only_initial_file(txt_file_path)

        await message.answer(f'🔙 Getting back'
                             f'\n\nMain Menu', reply_markup=bt.start_bot_buttons)
        await state.set_state(UserFlow.file_type)

    else:
        await message.answer('❌ Error. Invalid symbols please use buttons below 🤨')
        await state.set_state(UserFlow.txt_converting_file_type)

########################################################################################################################
                                                          # The End


